# загрузка библиотек
from datetime import datetime
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.action_chains import ActionChains
from selenium.common.exceptions import NoSuchElementException
from selenium.webdriver.common.keys import Keys
from selenium.common.exceptions import ElementClickInterceptedException
import time


# Опции для работы с сайтом, установка максимального расширения браузера, чтобы влез весь слайд
chrome_options = Options()
chrome_options.add_argument('--window-size=1920,1280')
chrome_options.add_argument('--lang=RU')
driver = webdriver.Chrome('\\\\pgk.rzd\\workgroups\\БД ЛП\\ЛПА (отдел аналитической работы)\\6. НАРАБОТКИ PYTHON\\Инструментарий Python\\Chromedriver\\chromedriver.exe', options=chrome_options)

# Опции для правого клика
action = ActionChains(driver)


# Объявление функции для снимка экрана и его обрезки.
'''Определяем директорию для сохранения слайда, делаем скрин экрана и присваиваем имя, далее открываем изображение,
обрезаем под новые параметры и сохраняем в ту же директорию с новым именем  '''
def screenie(i):
    dirscreen = r'\\pgk.rzd\workgroups\БД ЛП\ЛПА (отдел аналитической работы)\6. НАРАБОТКИ PYTHON\Роботы\Slide_maker_KR_Vostok\for_slide\main'
    driver.save_screenshot(dirscreen + str(i) + '.png')
    im = Image.open(dirscreen + str(i) + '.png')
    cropped = im.crop((10, 272, 1885, 800))                     # задаем размер обрезаемого изображения
    cropped.save(dirscreen + 'Cropped_' + str(i) + '.png')


# СЛАЙД
# Открываем браузер, выбираем фильтр РПС на КР, так как лист показывается только с этим фильтром, далее ждем 40 секунд, чтобы погрузилась страница
driver.get("https://bi.pgk.ru/self/sense/app/9ac4efd3-e568-4dbd-9335-7e8304328cf3/sheet/afe0b838-f241-4acc-aa48-ee9dac9694cf/state/analysis")
time.sleep(30)


# Вызываем функцию снимка экрана и его обрезаем
screenie(1)


# Закрытие браузера
driver.quit()


# Определение пустого шаблона презентации 16х9, задаем шаблон №6
dirslide = r'\\pgk.rzd\workgroups\БД ЛП\ЛПА (отдел аналитической работы)\6. НАРАБОТКИ PYTHON\Роботы\Slide_maker_KR_Vostok\for_slide\main'
prs = Presentation(dirslide + 'Sample.pptx')

# задаем параметры размера картинки, которая будет вставлена в слайд и выбираем номер макета из пустого слайда. Размер картинки определяется в дюймах
# если не задать размеры высоты и ширины, то размеры картинки вставляются исходными (конвертировал для удобства здесь https://www.unitconverters.net/typography/inch-to-pixel-x.htm)
left = Inches(0)
top = Inches(2)
height = Inches(4)
width = Inches(13.2)
blank_slide_layout = prs.slide_layouts[6]


# Формирование цикла для последовательного создания слайдов и вставки снимков экрана
for i in range(1, 2):
    # Добавление пустого слайда
    slide = prs.slides.add_slide(blank_slide_layout)

    # Вставка изображения на слайд
    img_path = dirslide + 'Cropped_' + str(i) + '.png'
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)


# Сохранение презентации
prs.save(dirslide + 'Slide.pptx')

# Формирование переменной с текущей датой
date_name = datetime.today()
date_format = date_name.strftime('%d.%m.%Y')

# Копирование файла с презентацией в другую папку Архив с присвоением даты
src = Path(dirslide + 'Slide.pptx')
dest = Path('//pgk.rzd/workgroups/БД ЛП/ЛПА (отдел аналитической работы)/6. НАРАБОТКИ PYTHON/Роботы/Slide_maker_KR_Vostok/slides_archive_new/Slide_' + str(date_format) + '.pptx')
dest.write_bytes(src.read_bytes())

# выход из скрипта
print('Slide_maker_KR_Vostok.py - COMPLETE', '\n', '___________________________________________________________', '\n')
